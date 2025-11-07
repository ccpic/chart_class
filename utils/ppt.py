from __future__ import annotations
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.shapes.autoshape import Shape
from pptx.shapes.picture import Picture
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.slide import SlideLayout, Slide
from pptx.table import Table
from typing import List, Tuple, Union, Optional
import math
import inspect
import pandas as pd

try:
    from typing import Literal
except ImportError:
    from typing_extensions import Literal


class Loc:
    def __init__(self, left: Union[Inches, Cm, int], top: Union[Inches, Cm, int]):
        self.left = left
        self.top = top
        """_初始化一个Loc类，定义一个ppt当中对象的位置
        
        Args:
            left(Union[Inches, Cm, int]): x轴坐标
            top(Union[Inches, Cm, int]): y轴坐标
        """

    def __repr__(self):
        return f"Loc({self.left}, {self.top})"

    def __iter__(self):
        yield self.left
        yield self.top

    def __add__(
        self,
        other: Union[
            Tuple[Union[Inches, Cm, int], Union[Inches, Cm, int]],
            List[Union[Inches, Cm, int], Union[Inches, Cm, int]],
        ],
    ) -> Loc:
        """与另一个二元元祖或列表相加，得到新位置

        Example:
            >>> loc1 = Loc(1, 2)
            >>> loc2 = Loc(3, 4)
            >>> loc3 = loc1 + loc2
            >>> loc3
            Loc(4, 6)
            >>> loc4 = loc1 + (5, 6)
            >>> loc4
            Loc(6, 8)

        Args:
            other (Union[
            Tuple[Union[Inches, Cm, int], Union[Inches, Cm, int]],
            List[Union[Inches, Cm, int], Union[Inches, Cm, int]],
            ]): 要相加的二元元祖或列表

        Raises:
            TypeError: 类型错误

        Returns:
            Loc: 返回新位置
        """

        if isinstance(other, (tuple, list)) and len(other) == 2:
            return Loc(self.left + other[0], self.top + other[1])
        else:
            raise TypeError("不支持的类型，只支持二元元祖或列表。")

    def __sub__(self, other):
        """与另一个二元元祖或列表相减，得到新位置

        Example:
            >>> loc1 = Loc(1, 2)
            >>> loc2 = Loc(3, 4)
            >>> loc3 = loc1 - loc2
            >>> loc3
            Loc(-2, -2)
            >>> loc4 = loc1 - (5, 6)
            >>> loc4
            Loc(-4, -4)

        Args:
            other (Union[
            Tuple[Union[Inches, Cm, int], Union[Inches, Cm, int]],
            List[Union[Inches, Cm, int], Union[Inches, Cm, int]],
            ]): 要相加的二元元祖或列表

        Raises:
            TypeError: 类型错误

        Returns:
            Loc: 返回新位置
        """
        if isinstance(other, (tuple, list)) and len(other) == 2:
            return Loc(self.left - other[0], self.top - other[1])
        else:
            raise TypeError("不支持的类型，只支持二元元祖或列表。")


class AnchorLoc(Loc):
    def __init__(
        self,
        left: Union[Inches, Cm, int],
        top: Union[Inches, Cm, int],
        shape_width: Union[Inches, Cm, int],
        shape_height: Union[Inches, Cm, int],
        anchor: Literal[
            "center",
            "top_left",
            "left_top",
            "top_right",
            "right_top",
            "top_mid",
            "mid_top",
            "top",
            "bottom_right",
            "right_bottom",
            "bottom_left",
            "left_bottom",
            "bottom_mid",
            "mid_bottom",
            "bottom",
            "mid_right",
            "right_mid",
            "right",
            "mid_left",
            "right_left",
            "left",
        ] = "center",
    ):
        """_初始化一个AnchorLoc类，定义一个ppt当中对象的位置，并根据锚点进行调整

        Args:
            left(Union[Inches, Cm, int]): x轴坐标
            top(Union[Inches, Cm, int]): y轴坐标
            shape_width(Union[Inches, Cm, int]): 对象的宽度
            shape_height(Union[Inches, Cm, int]): 对象的高度
            anchor (Literal[
                "center",
                "top_left", "left_top",
                "top_right", "right_top",
                "top_mid", "mid_top", "top",
                "bottom_right", "right_bottom",
                "bottom_left", "left_bottom",
                "bottom_mid", "mid_bottom", "bottom",
                "mid_right", "right_mid", "right",
                "mid_left", "right_left", "left",
            ]): 锚点. Defaults to "center".

        """
        super().__init__(left, top)  # 调用父类的构造函数，初始化 left 和 top
        self.shape_width = shape_width
        self.shape_height = shape_height
        self.anchor = anchor

    def __repr__(self):
        return f"AnchorLoc(left={self.left}, top={self.top}, shape_width={self.shape_width}, shape_height={self.shape_height}), anchor={self.anchor})"

    @property
    def loc(self) -> Loc:
        """根据形状的宽高和锚点计算调整后的坐标

        Returns:
            Loc: 返回调整后的Loc对象
        """
        if self.anchor == "center":
            left = self.left - self.shape_width / 2
            top = self.top - self.shape_height / 2
        elif self.anchor in ["top_right", "right_top"]:
            left = self.left - self.shape_width
            top = self.top
        elif self.anchor in ["top_left", "left_top"]:
            left = self.left
            top = self.top
        elif self.anchor in ["top_mid", "mid_top", "top"]:
            left = self.left - self.shape_width / 2
            top = self.top
        elif self.anchor in ["bottom_right", "right_bottom"]:
            left = self.left - self.shape_width
            top = self.top - self.shape_height
        elif self.anchor in ["bottom_left", "left_bottom"]:
            left = self.left
            top = self.top - self.shape_height
        elif self.anchor in ["bottom_mid", "mid_bottom", "bottom"]:
            left = self.left - self.shape_width / 2
            top = self.top - self.shape_height
        elif self.anchor in ["mid_right", "right_mid", "right"]:
            left = self.left - self.shape_width
            top = self.top - self.shape_height / 2
        elif self.anchor in ["mid_left", "right_left", "left"]:
            left = self.left
            top = self.top - self.shape_height / 2

        return Loc(int(left), int(top))  # pptx很多参数如shape的top和left只接受整数


class Section:
    def __init__(
        self,
        left: Union[Inches, Cm, int],
        top: Union[Inches, Cm, int],
        width: Union[Inches, Cm, int],
        height: Union[Inches, Cm, int],
    ) -> None:
        """初始化Section类，定义一个矩形范围

        Parameters
        ----------
        left: Union[Inches, Cm, int]
            左边距
        right: Union[Inches, Cm, int]
            右边距
        width: Union[Inches, Cm, int]
            宽度
        height: Union[Inches, Cm, int]
            高度
        """
        self.left = left
        self.top = top
        self.width = width
        self.height = height
        self.right = self.left + self.width
        self.bottom = self.top + self.height

    def __repr__(self):
        return f"Section({self.left}, {self.top}, {self.width}, {self.height})"

    @property
    def left_top(self) -> Loc:
        """Section左上角坐标

        Returns:
            Loc: 返回一个(左边距, 上边距)的tuple
        """
        return Loc(self.left, self.top)

    @property
    def mid_top(self) -> Loc:
        """Section上边居中位置坐标

        Returns:
            Loc: 返回一个(左边距, 上边距)的tuple
        """
        return Loc(self.left + (self.width / 2), self.top)

    @property
    def right_top(self) -> Loc:
        """Section右上角坐标

        Returns:
            Loc: 返回一个(左边距, 上边距)的tuple
        """
        return Loc(self.right, self.top)

    @property
    def left_mid(self) -> Loc:
        """Section左边居中位置坐标

        Returns:
            Loc: 返回一个(左边距, 上边距)的tuple
        """
        return Loc(self.left, self.top + (self.height / 2))

    @property
    def center(self) -> Loc:
        """Section整体居中位置坐标

        Returns:
            Loc: 返回一个(左边距, 上边距)的tuple
        """
        return Loc(self.left + (self.width / 2), self.top + (self.height / 2))

    @property
    def right_mid(self) -> Loc:
        """Section右边居中位置坐标

        Returns:
            Loc: 返回一个(左边距, 上边距)的tuple
        """
        return Loc(self.right, self.top + (self.height / 2))

    @property
    def left_bottom(self) -> Loc:
        """Section左下角坐标

        Returns:
            Loc: 返回一个(左边距, 上边距)的tuple
        """
        return Loc(self.left, self.bottom)

    @property
    def mid_bottom(self) -> Loc:
        """Section下边居中位置坐标

        Returns:
            Loc: 返回一个(左边距, 上边距)的tuple
        """
        return Loc(self.left + (self.width / 2), self.bottom)

    @property
    def right_bottom(self) -> Loc:
        """Section右下角坐标

        Returns:
            Loc: 返回一个(左边距, 上边距)的tuple
        """
        return Loc(self.right, self.bottom)

    def fraction(
        self,
        dimension: Literal["width", "height"],
        frac_n: int,
        index: Optional[int] = 1,
    ) -> Section:
        """宽度/高度等分并返回指定部分

        Args:
            dimension(Literal["width", "height"]): 等分宽度还是高度
            frac_n(int): 宽度等分成几份
            index(Optional[int]), optional: 返回等分后的第几份，不能大于frac_n. Defaults to 1.

        Returns:
            Section: 返回自身实例
        """

        if index > frac_n:
            raise ValueError(f"参数'index'的值{index}不能大于参数'frac_n'的值{frac_n}")

        if dimension == "width":
            width_new = self.width / frac_n
            left_new = self.left + width_new * (index - 1)

            self.width = width_new
            self.left = left_new
        elif dimension == "height":
            height_new = self.height / frac_n
            top_new = self.top + height_new * (index - 1)

            self.height = height_new
            self.top = top_new

        return self

    def apply_margin(
        self,
        dimension: Literal["width", "height"] = "width",
        margin: Union[Inches, Cm, int] = Cm(0.64),
    ) -> Section:
        """根据边距调整区域

        Args:
            dimension(Literal["width", "height"]): 调整宽度边距还是高度边距. Defaults to "width".
            margin(Union[Inches, Cm, int]): 宽度等分成几份: 边距值，默认0.64厘米. Defaults to Cm(0.64).

        Returns:
            Section: 返回自身实例
        """
        if dimension == "width":
            self.width = self.width - margin * 2
            self.left = self.left + margin
        if dimension == "height":
            self.height = self.height - margin * 2
            self.top = self.top + margin

        return self

    """alias"""
    top_left = left_top
    top_mid = mid_top
    top_right = right_top
    mid_left = left_mid
    mid = center
    mid_right = right_mid
    bottom_left = left_bottom
    bottom_mid = mid_bottom
    bottom_right = right_bottom


def is_light_color(color: Optional[Union[RGBColor, str]]) -> bool:
    """判断一个颜色是否是浅色

    Args:
        color (Optional[Union[RGBColor, str]]): 一个颜色，可以是pptx包的RGBColor类型，也可以是十六进制颜色字符串

    Returns:
        bool: 返回一个布尔值
    """
    if color is None:
        return True
    else:
        rgb_color = (
            color if isinstance(color, RGBColor) else RGBColor.from_string(color)
        )
        r, g, b = rgb_color
        hsp = math.sqrt(0.299 * (r * r) + 0.587 * (g * g) + 0.114 * (b * b))

        if hsp > 127.5:
            return True
        else:
            return False


class SlideContent:
    def __init__(
        self,
        prs: Presentation,
        slide: Slide,
        header_height: Cm = Cm(2.72),
        body_top: Cm = Cm(2.91),
        body_height: Cm = Cm(13.81),
        footer_top: Cm = Cm(17.36),
        footer_height: Cm = Cm(1.71),
    ) -> None:
        self.prs = prs
        self.slide = slide
        self.header_height = header_height
        self.body_top = body_top
        self.body_height = body_height
        self.footer_top = footer_top
        self.footer_height = footer_height

    @property
    def header(self) -> Section:
        """Slide的页眉区域

        Returns:
            Section: 返回一个Section对象
        """
        return Section(
            left=Cm(0), top=Cm(0), width=self.prs.slide_width, height=self.header_height
        )

    @property
    def body(self) -> Section:
        """Slide的页面主体区域

        Returns:
            Section: 返回一个Section对象
        """
        return Section(
            left=Cm(0),
            top=self.body_top,
            width=self.prs.slide_width,
            height=self.body_height,
        )

    @property
    def footer(self) -> Section:
        """Slide的页脚区域

        Returns:
            Section: 返回一个Section对象
        """
        return Section(
            left=Cm(5.6),
            top=self.footer_top,
            width=self.prs.slide_width,
            height=self.footer_height,
        )

    def add_text(
        self,
        text: str,
        width: Union[Inches, Cm, int],
        height: Union[Inches, Cm, int],
        loc: Loc,
        anchor: Literal[
            "center",
            "top_left",
            "left_top",
            "top_right",
            "right_top",
            "top_mid",
            "mid_top",
            "top",
            "bottom_right",
            "right_bottom",
            "bottom_left",
            "left_bottom",
            "bottom_mid",
            "mid_bottom",
            "bottom",
            "mid_right",
            "right_mid",
            "right",
            "mid_left",
            "right_left",
            "left",
        ] = "center",
        shape_type: MSO_SHAPE_TYPE = MSO_SHAPE.RECTANGLE,
        fill_color: Optional[Union[RGBColor, str]] = None,
        border_color: Optional[Union[RGBColor, str]] = None,
        font_color: Optional[Union[RGBColor, str]] = None,
        font_size: Optional[Union[Pt, int]] = Pt(14),
        font_bold=False,
        font_italic=False,
        text_wrap: bool = False,
        *args,
        **kwargs,
    ) -> Shape:
        """_summary_

        Args:
            text (str): 文本内容
            width (Union[Inches, Cm, int]): 宽度
            height (Union[Inches, Cm, int]): 高度
            loc (Loc - left, top): _description_
            anchor (Literal[
                "center",
                "top_left", "left_top",
                "top_right", "right_top",
                "top_mid", "mid_top", "top",
                "bottom_right", "right_bottom",
                "bottom_left", "left_bottom",
                "bottom_mid", "mid_bottom", "bottom",
                "mid_right", "right_mid", "right",
                "mid_left", "right_left", "left",
            ]): 锚点. Defaults to "center".
            fill_color (Optional[Union[RGBColor, str]], optional): 填充颜色. Defaults to None.
            border_color (Optional[Union[RGBColor, str]], optional): 边框颜色. Defaults to None.
            font_color (Optional[Union[RGBColor, str]], optional): 字体颜色. Defaults to None.
            font_size (Optional[Union[Pt, int]], optional): 文字大小. Defaults to Pt(14).
            font_bold (bool, optional): 文字是否粗体. Defaults to False.
            font_italic (bool, optional): 文字是否斜体. Defaults to False.
            text_wrap (bool, optional): 文字是否换行. Defaults to False.

        Returns:
            Shape: 带有文本的形状
        """
        left, top = AnchorLoc(loc.left, loc.top, width, height, anchor=anchor).loc
        shape = self.slide.shapes.add_shape(
            shape_type,
            left=left,
            top=top,
            width=width,
            height=height,
        )

        # 形状填充
        fill = shape.fill
        if fill_color is None:
            fill.background()
        else:
            fill.solid()
            fill.fore_color.rgb = (
                fill_color
                if isinstance(fill_color, RGBColor)
                else RGBColor.from_string(fill_color)
            )

        # 边框
        line = shape.line
        if border_color is None:
            line.fill.background()
        else:
            line.color.rgb = (
                border_color
                if isinstance(border_color, RGBColor)
                else RGBColor.from_string(border_color)
            )

        # 形状文字
        text_frame = shape.text_frame
        text_frame.word_wrap = text_wrap
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text

        # 设置字体
        font = run.font
        font.name = "微软雅黑"
        font.size = font_size if isinstance(font_size, Pt) else Pt(font_size)
        font.bold = font_bold
        font.italic = font_italic
        if font_color is None:
            if is_light_color(fill_color):
                font.color.rgb = RGBColor(0, 0, 0)
            else:
                font.color.rgb = RGBColor(255, 255, 255)
        else:
            font.color.rgb = (
                font_color
                if isinstance(font_color, RGBColor)
                else RGBColor.from_string(font_color)
            )

        return shape

    def add_image(
        self,
        image_file: str,
        width: Optional[Union[Inches, Cm, int]] = None,
        height: Optional[Union[Inches, Cm, int]] = None,
        top: Optional[Union[Inches, Cm, int]] = None,
        left: Optional[Union[Inches, Cm, int]] = None,
        loc: Optional[Loc] = None,
        anchor: Optional[
            Literal[
                "center",
                "top_left",
                "left_top",
                "top_right",
                "right_top",
                "top_mid",
                "mid_top",
                "top",
                "bottom_right",
                "right_bottom",
                "bottom_left",
                "left_bottom",
                "bottom_mid",
                "mid_bottom",
                "bottom",
                "mid_right",
                "right_mid",
                "right",
                "mid_left",
                "right_left",
                "left",
            ]
        ] = "center",
    ) -> Picture:
        # 先插入图片，方便获取图片长宽（插入图片时如不同时指定宽和高，pptx高会自动缩放尺寸），再调整位置
        image = self.slide.shapes.add_picture(
            image_file=image_file,
            left=0 if left is None else left,
            top=0 if top is None else top,
            width=width,
            height=height,
        )

        if loc is not None:
            left, top = AnchorLoc(
                loc.left, loc.top, image.width, image.height, anchor=anchor
            ).loc
            image.left = left
            image.top = top

        return image

    def set_title(
        self, title: str, font_size: Optional[Union[Pt, Inches, Cm]] = None
    ) -> Shape:
        """设置幻灯片标题

        Args:
            title (str): 标题内容
            font_size (Optional[Union[Pt, Inches, Cm]], optional): 标题字体大小. Defaults to None.

        Returns:
            Shape: 带有标题的形状
        """

        title_shape = self.slide.shapes.title
        title_shape.text = title

        # 设置字体大小
        if font_size is not None:
            text_frame = title_shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = font_size

        return title_shape

    def add_table(
        self,
        df: pd.DataFrame,
        width: Optional[Union[Inches, Cm, int]] = None,
        height: Optional[Union[Inches, Cm, int]] = None,
        top: Optional[Union[Inches, Cm, int]] = None,
        left: Optional[Union[Inches, Cm, int]] = None,
        loc: Optional[Loc] = None,
        anchor: Optional[
            Literal[
                "center",
                "top_left",
                "left_top",
                "top_right",
                "right_top",
                "top_mid",
                "mid_top",
                "top",
                "bottom_right",
                "right_bottom",
                "bottom_left",
                "left_bottom",
                "bottom_mid",
                "mid_bottom",
                "bottom",
                "mid_right",
                "right_mid",
                "right",
                "mid_left",
                "right_left",
                "left",
            ]
        ] = "center",
        **kwargs,
    ) -> Table:
        """将pandas DataFrame添加为PowerPoint表格

        Args:
            df (pd.DataFrame): 要转换的DataFrame
            width (Optional[Union[Inches, Cm, int]], optional): 表格宽度. Defaults to None.
            height (Optional[Union[Inches, Cm, int]], optional): 表格高度. Defaults to None.
            top (Optional[Union[Inches, Cm, int]], optional): 表格顶部位置. Defaults to None.
            left (Optional[Union[Inches, Cm, int]], optional): 表格左侧位置. Defaults to None.
            loc (Optional[Loc], optional): 表格位置. Defaults to None.
            anchor (Optional[Literal[...]], optional): 锚点位置. Defaults to "center".

        Keyword Args:
            show_index (bool): 是否显示DataFrame的索引，默认为True
            header_bg_color (Union[RGBColor, str]): 表头背景颜色
            header_font_color (Union[RGBColor, str]): 表头字体颜色
            header_font_size (Union[Pt, int]): 表头字体大小
            header_font_bold (bool): 表头是否粗体
            body_bg_color (Union[RGBColor, str]): 表格主体背景颜色
            body_font_color (Union[RGBColor, str]): 表格主体字体颜色
            body_font_size (Union[Pt, int]): 表格主体字体大小
            border_color (Union[RGBColor, str]): 边框颜色
            border_width (Union[Pt, int]): 边框宽度
            row_height (Union[Inches, Cm, int]): 行高
            col_width (Union[Inches, Cm, int]): 列宽
            first_row_bold (bool): 第一行是否加粗
            zebra_stripes (bool): 是否启用斑马纹
            zebra_color (Union[RGBColor, str]): 斑马纹颜色

        Returns:
            Table: PowerPoint表格对象
        """
        # 获取DataFrame的行数和列数
        rows, cols = df.shape

        # 检查是否需要显示索引
        show_index = kwargs.get("show_index", True)

        # 如果显示索引，列数需要+1
        if show_index:
            cols += 1

        # 设置默认尺寸
        if width is None:
            width = Cm(15)  # 默认宽度15厘米
        if height is None:
            height = Cm(rows * 0.8 + 1)  # 根据行数计算默认高度

        # 计算默认位置
        if loc is None and top is None and left is None:
            loc = self.body.center
            anchor = "center"

        # 创建表格
        table = self.slide.shapes.add_table(
            rows=rows + 1,  # +1 for header
            cols=cols,
            left=0 if left is None else left,
            top=0 if top is None else top,
            width=width,
            height=height,
        ).table

        # 设置表头
        col_offset = 0
        if show_index:
            # 第一列显示索引名称
            index_name = df.index.name if df.index.name else "索引"
            cell = table.cell(0, 0)
            cell.text = str(index_name)
            col_offset = 1

        # 设置数据列的表头
        for col_idx, col_name in enumerate(df.columns):
            cell = table.cell(0, col_idx + col_offset)
            cell.text = str(col_name)

        # 填充数据
        for row_idx in range(rows):
            if show_index:
                # 第一列显示索引值
                cell = table.cell(row_idx + 1, 0)
                cell.text = str(df.index[row_idx])

            # 填充数据列
            for col_idx in range(len(df.columns)):
                cell = table.cell(row_idx + 1, col_idx + col_offset)
                cell.text = str(df.iloc[row_idx, col_idx])

        # 应用样式设置
        self._apply_table_styling(table, **kwargs)

        # 调整位置（如果需要）
        if loc is not None:
            left, top = AnchorLoc(loc.left, loc.top, width, height, anchor=anchor).loc
            # 获取表格的图形对象来设置位置
            # 通过表格的图形对象来设置位置
            table_shape = table._graphic_frame
            table_shape.left = left
            table_shape.top = top

        return table

    def _apply_table_styling(self, table: Table, **kwargs):
        """应用表格样式设置

        Args:
            table (Table): PowerPoint表格对象
            **kwargs: 样式参数
        """
        # 获取样式参数
        header_bg_color = kwargs.get("header_bg_color", RGBColor(68, 114, 196))
        header_font_color = kwargs.get("header_font_color", RGBColor(255, 255, 255))
        header_font_size = kwargs.get("header_font_size", Pt(12))
        header_font_bold = kwargs.get("header_font_bold", True)

        body_bg_color = kwargs.get("body_bg_color", RGBColor(255, 255, 255))
        body_font_color = kwargs.get("body_font_color", RGBColor(0, 0, 0))
        body_font_size = kwargs.get("body_font_size", Pt(10))

        border_color = kwargs.get("border_color", RGBColor(217, 217, 217))
        border_width = kwargs.get("border_width", Pt(1))

        row_height = kwargs.get("row_height", Cm(0.8))
        col_width = kwargs.get("col_width", None)

        first_row_bold = kwargs.get("first_row_bold", True)
        zebra_stripes = kwargs.get("zebra_stripes", False)
        zebra_color = kwargs.get("zebra_color", RGBColor(242, 242, 242))

        # 设置表头样式
        for col_idx in range(len(table.columns)):
            cell = table.cell(0, col_idx)
            # 背景色
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                header_bg_color
                if isinstance(header_bg_color, RGBColor)
                else RGBColor.from_string(header_bg_color)
            )
            # 字体
            paragraph = cell.text_frame.paragraphs[0]
            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
            font = run.font
            font.size = (
                header_font_size
                if isinstance(header_font_size, Pt)
                else Pt(header_font_size)
            )
            font.bold = header_font_bold
            font.color.rgb = (
                header_font_color
                if isinstance(header_font_color, RGBColor)
                else RGBColor.from_string(header_font_color)
            )

        # 设置表格主体样式
        for row_idx in range(1, len(table.rows)):
            for col_idx in range(len(table.columns)):
                cell = table.cell(row_idx, col_idx)

                # 背景色（斑马纹）
                if zebra_stripes and row_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = (
                        zebra_color
                        if isinstance(zebra_color, RGBColor)
                        else RGBColor.from_string(zebra_color)
                    )
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = (
                        body_bg_color
                        if isinstance(body_bg_color, RGBColor)
                        else RGBColor.from_string(body_bg_color)
                    )

                # 字体
                paragraph = cell.text_frame.paragraphs[0]
                run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                font = run.font
                font.size = (
                    body_font_size
                    if isinstance(body_font_size, Pt)
                    else Pt(body_font_size)
                )
                font.bold = first_row_bold and row_idx == 1
                font.color.rgb = (
                    body_font_color
                    if isinstance(body_font_color, RGBColor)
                    else RGBColor.from_string(body_font_color)
                )

        # 设置边框 - 通过表格的边框属性来设置
        # 注意：pptx库中单元格没有直接的border属性，需要通过表格的边框设置
        # 这里我们设置表格的整体边框样式
        for row in table.rows:
            for cell in row.cells:
                # 设置单元格填充
                if not cell.fill.fore_color.rgb:  # 如果没有设置填充色，设置默认填充
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = (
                        body_bg_color
                        if isinstance(body_bg_color, RGBColor)
                        else RGBColor.from_string(body_bg_color)
                    )

        # 设置行高和列宽
        for row in table.rows:
            row.height = row_height if isinstance(row_height, int) else int(row_height)

        if col_width is not None:
            for col in table.columns:
                col.width = col_width if isinstance(col_width, int) else int(col_width)


class PPT:
    def __init__(self, template_path: str) -> None:
        """初始化一个PPT类

        Args:
            template_path (str): 模板文件路径
            save_path (str): 保存路径，默认为和模板路径同一文件夹下的presentation.pptx文件. Defaults to None.
        """
        self.template_path = template_path
        self.prs = Presentation(template_path)

    @property
    def parent_slide(self) -> SlideLayout:
        """内容页幻灯片母版

        Returns
        -------
        SlideLayout
            返回该ppt的内容页幻灯片母版对象
        """
        slide = self.prs.slide_layouts[0]
        return slide

    def add_content_slide(
        self,
        layout_style: int = 0,
        header_height: Cm = Cm(2.72),
        body_top: Cm = Cm(2.91),
        body_height: Cm = Cm(13.81),
        footer_top: Cm = Cm(17.36),
        footer_height: Cm = Cm(1.71),
    ) -> SlideContent:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_style])
        print(f"已添加第{self.prs.slides.index(slide)}页.")
        content = SlideContent(
            self.prs,
            slide,
            header_height,
            body_top,
            body_height,
            footer_top,
            footer_height,
        )
        return content

    def remove_slide(self, index: int) -> None:
        """删除指定索引的幻灯片

        Args:
            index (int): 幻灯片索引
        """

        # 检查索引是否有效
        if index < 0 or index >= len(self.prs.slides):
            raise ValueError("幻灯片索引超出范围")

        rId = self.prs.slides._sldIdLst[index].rId
        self.prs.part.drop_rel(rId)
        del self.prs.slides._sldIdLst[index]
        print(f"已删除第{index+1}页.")

    def save(self, save_path: str) -> None:
        # 获取调用当前脚本的脚本路径
        calling_script_dir = os.path.dirname(
            os.path.abspath(inspect.stack()[1].filename)
        )
        ppt_dir = f"{calling_script_dir}"

        # 保存
        if os.path.exists(ppt_dir) is False:
            os.makedirs(ppt_dir)

        self.prs.save(save_path)
        print(f"{save_path}已保存.")


if __name__ == "__main__":
    p = PPT("template.pptx")
    c = p.add_content_slide()
    c.set_title("测试测试")
    c.add_text(
        "测试",
        width=Cm(4),
        height=Cm(0.5),
        loc=c.body.top_right,
        anchor="top_right",
        fill_color="F0CB46",
    )
    c.add_image(
        "D:\PyProjects\chart_class\plots\柱状图.png",
        width=c.body.width * 0.8,
        height=None,
        loc=c.body.center,
    )
    p.save()
