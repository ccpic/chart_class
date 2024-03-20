from __future__ import annotations
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.shapes.autoshape import Shape
from pptx.shapes.picture import Picture
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.slide import SlideLayout, Slide
from typing import List, Tuple, Union, Optional
import math

try:
    from typing import Literal
except ImportError:
    from typing_extensions import Literal


class Loc:
    def __init__(self, left: Union[Inches, Cm, int], top: Union[Inches, Cm, int]):
        self.left = left
        self.top = top
        """_初始化一个Loc类，定义一个ppt当中对象的位置
        
        Args:
            left(Union[Inches, Cm, int]): x轴坐标
            top(Union[Inches, Cm, int]): y轴坐标
        """

    def __repr__(self):
        return f"Loc({self.left}, {self.top})"

    def __iter__(self):
        yield self.left
        yield self.top

    def __add__(
        self,
        other: Union[
            Tuple[Union[Inches, Cm, int], Union[Inches, Cm, int]],
            List[Union[Inches, Cm, int], Union[Inches, Cm, int]],
        ],
    ) -> Loc:
        """与另一个二元元祖或列表相加，得到新位置

        Example:
            >>> loc1 = Loc(1, 2)
            >>> loc2 = Loc(3, 4)
            >>> loc3 = loc1 + loc2
            >>> loc3
            Loc(4, 6)
            >>> loc4 = loc1 + (5, 6)
            >>> loc4
            Loc(6, 8)

        Args:
            other (Union[
            Tuple[Union[Inches, Cm, int], Union[Inches, Cm, int]],
            List[Union[Inches, Cm, int], Union[Inches, Cm, int]],
            ]): 要相加的二元元祖或列表

        Raises:
            TypeError: 类型错误

        Returns:
            Loc: 返回新位置
        """

        if isinstance(other, (tuple, list)) and len(other) == 2:
            return Loc(self.left + other[0], self.top + other[1])
        else:
            raise TypeError("不支持的类型，只支持二元元祖或列表。")

    def __sub__(self, other):
        """与另一个二元元祖或列表相减，得到新位置

        Example:
            >>> loc1 = Loc(1, 2)
            >>> loc2 = Loc(3, 4)
            >>> loc3 = loc1 - loc2
            >>> loc3
            Loc(-2, -2)
            >>> loc4 = loc1 - (5, 6)
            >>> loc4
            Loc(-4, -4)

        Args:
            other (Union[
            Tuple[Union[Inches, Cm, int], Union[Inches, Cm, int]],
            List[Union[Inches, Cm, int], Union[Inches, Cm, int]],
            ]): 要相加的二元元祖或列表

        Raises:
            TypeError: 类型错误

        Returns:
            Loc: 返回新位置
        """
        if isinstance(other, (tuple, list)) and len(other) == 2:
            return Loc(self.left - other[0], self.top - other[1])
        else:
            raise TypeError("不支持的类型，只支持二元元祖或列表。")


class AnchorLoc(Loc):
    def __init__(
        self,
        left: Union[Inches, Cm, int],
        top: Union[Inches, Cm, int],
        shape_width: Union[Inches, Cm, int],
        shape_height: Union[Inches, Cm, int],
        anchor: Literal[
            "center",
            "top_left",
            "left_top",
            "top_right",
            "right_top",
            "top_mid",
            "mid_top",
            "top",
            "bottom_right",
            "right_bottom",
            "bottom_left",
            "left_bottom",
            "bottom_mid",
            "mid_bottom",
            "bottom",
            "mid_right",
            "right_mid",
            "right",
            "mid_left",
            "right_left",
            "left",
        ] = "center",
    ):
        """_初始化一个AnchorLoc类，定义一个ppt当中对象的位置，并根据锚点进行调整

        Args:
            left(Union[Inches, Cm, int]): x轴坐标
            top(Union[Inches, Cm, int]): y轴坐标
            shape_width(Union[Inches, Cm, int]): 对象的宽度
            shape_height(Union[Inches, Cm, int]): 对象的高度
            anchor (Literal[
                "center",
                "top_left", "left_top",
                "top_right", "right_top",
                "top_mid", "mid_top", "top",
                "bottom_right", "right_bottom",
                "bottom_left", "left_bottom",
                "bottom_mid", "mid_bottom", "bottom",
                "mid_right", "right_mid", "right",
                "mid_left", "right_left", "left",
            ]): 锚点. Defaults to "center".

        """
        super().__init__(left, top)  # 调用父类的构造函数，初始化 left 和 top
        self.shape_width = shape_width
        self.shape_height = shape_height
        self.anchor = anchor

    def __repr__(self):
        return f"AnchorLoc(left={self.left}, top={self.top}, shape_width={self.shape_width}, shape_height={self.shape_height}), anchor={self.anchor})"

    @property
    def loc(self) -> Loc:
        """根据形状的宽高和锚点计算调整后的坐标

        Returns:
            Loc: 返回调整后的Loc对象
        """
        if self.anchor == "center":
            left = self.left - self.shape_width / 2
            top = self.top - self.shape_height / 2
        elif self.anchor in ["top_right", "right_top"]:
            left = self.left - self.shape_width
            top = self.top
        elif self.anchor in ["top_left", "left_top"]:
            left = self.left
            top = self.top
        elif self.anchor in ["top_mid", "mid_top", "top"]:
            left = self.left - self.shape_width / 2
            top = self.top
        elif self.anchor in ["bottom_right", "right_bottom"]:
            left = self.left - self.shape_width
            top = self.top - self.shape_height
        elif self.anchor in ["bottom_left", "left_bottom"]:
            left = self.left
            top = self.top - self.shape_height
        elif self.anchor in ["bottom_mid", "mid_bottom", "bottom"]:
            left = self.left - self.shape_width / 2
            top = self.top - self.shape_height
        elif self.anchor in ["mid_right", "right_mid", "right"]:
            left = self.left - self.shape_width
            top = self.top - self.shape_height / 2
        elif self.anchor in ["mid_left", "right_left", "left"]:
            left = self.left
            top = self.top - self.shape_height / 2

        return Loc(int(left), int(top))  # pptx很多参数如shape的top和left只接受整数


class Section:
    def __init__(
        self,
        left: Union[Inches, Cm, int],
        top: Union[Inches, Cm, int],
        width: Union[Inches, Cm, int],
        height: Union[Inches, Cm, int],
    ) -> None:
        """初始化Section类，定义一个矩形范围

        Parameters
        ----------
        left: Union[Inches, Cm, int]
            左边距
        right: Union[Inches, Cm, int]
            右边距
        width: Union[Inches, Cm, int]
            宽度
        height: Union[Inches, Cm, int]
            高度
        """
        self.left = left
        self.top = top
        self.width = width
        self.height = height
        self.right = self.left + self.width
        self.bottom = self.top + self.height

    def __repr__(self):
        return f"Section({self.left}, {self.top}, {self.width}, {self.height})"

    @property
    def left_top(self) -> Loc:
        """Section左上角坐标

        Returns:
            Loc: 返回一个(左边距, 上边距)的tuple
        """
        return Loc(self.left, self.top)

    @property
    def mid_top(self) -> Loc:
        """Section上边居中位置坐标

        Returns:
            Loc: 返回一个(左边距, 上边距)的tuple
        """
        return Loc(self.left + (self.width / 2), self.top)

    @property
    def right_top(self) -> Loc:
        """Section右上角坐标

        Returns:
            Loc: 返回一个(左边距, 上边距)的tuple
        """
        return Loc(self.right, self.top)

    @property
    def left_mid(self) -> Loc:
        """Section左边居中位置坐标

        Returns:
            Loc: 返回一个(左边距, 上边距)的tuple
        """
        return Loc(self.left, self.top + (self.height / 2))

    @property
    def center(self) -> Loc:
        """Section整体居中位置坐标

        Returns:
            Loc: 返回一个(左边距, 上边距)的tuple
        """
        return Loc(self.left + (self.width / 2), self.top + (self.height / 2))

    @property
    def right_mid(self) -> Loc:
        """Section右边居中位置坐标

        Returns:
            Loc: 返回一个(左边距, 上边距)的tuple
        """
        return Loc(self.right, self.top + (self.height / 2))

    @property
    def left_bottom(self) -> Loc:
        """Section左下角坐标

        Returns:
            Loc: 返回一个(左边距, 上边距)的tuple
        """
        return Loc(self.left, self.bottom)

    @property
    def mid_bottom(self) -> Loc:
        """Section下边居中位置坐标

        Returns:
            Loc: 返回一个(左边距, 上边距)的tuple
        """
        return Loc(self.left + (self.width / 2), self.bottom)

    @property
    def right_bottom(self) -> Loc:
        """Section右下角坐标

        Returns:
            Loc: 返回一个(左边距, 上边距)的tuple
        """
        return Loc(self.right, self.bottom)

    def fraction(
        self,
        dimension: Literal["width", "height"],
        frac_n: int,
        index: Optional[int] = 1,
    ) -> Section:
        """宽度/高度等分并返回指定部分

        Args:
            dimension(Literal["width", "height"]): 等分宽度还是高度
            frac_n(int): 宽度等分成几份
            index(Optional[int]), optional: 返回等分后的第几份，不能大于frac_n. Defaults to 1.

        Returns:
            Section: 返回自身实例
        """

        if index > frac_n:
            raise ValueError(f"参数'index'的值{index}不能大于参数'frac_n'的值{frac_n}")

        if dimension == "width":
            width_new = self.width / frac_n
            left_new = self.left + width_new * (index - 1)

            self.width = width_new
            self.left = left_new
        elif dimension == "height":
            height_new = self.height / frac_n
            top_new = self.top + height_new * (index - 1)

            self.height = height_new
            self.top = top_new

        return self

    def apply_margin(
        self,
        dimension: Literal["width", "height"] = "width",
        margin: Union[Inches, Cm, int] = Cm(0.64),
    ) -> Section:
        """根据边距调整区域

        Args:
            dimension(Literal["width", "height"]): 调整宽度边距还是高度边距. Defaults to "width".
            margin(Union[Inches, Cm, int]): 宽度等分成几份: 边距值，默认0.64厘米. Defaults to Cm(0.64).

        Returns:
            Section: 返回自身实例
        """
        if dimension == "width":
            self.width = self.width - margin * 2
            self.left = self.left + margin
        if dimension == "height":
            self.height = self.height - margin * 2
            self.top = self.top + margin

        return self

    """alias"""
    top_left = left_top
    top_mid = mid_top
    top_right = right_top
    mid_left = left_mid
    mid = center
    mid_right = right_mid
    bottom_left = left_bottom
    bottom_mid = mid_bottom
    bottom_right = right_bottom


def is_light_color(color: Optional[Union[RGBColor, str]]) -> bool:
    """判断一个颜色是否是浅色

    Args:
        color (Optional[Union[RGBColor, str]]): 一个颜色，可以是pptx包的RGBColor类型，也可以是十六进制颜色字符串

    Returns:
        bool: 返回一个布尔值
    """
    if color is None:
        return True
    else:
        rgb_color = (
            color if isinstance(color, RGBColor) else RGBColor.from_string(color)
        )
        r, g, b = rgb_color
        hsp = math.sqrt(0.299 * (r * r) + 0.587 * (g * g) + 0.114 * (b * b))

        if hsp > 127.5:
            return True
        else:
            return False


class SlideContent:
    def __init__(self, prs: Presentation, slide: Slide) -> None:
        self.prs = prs
        self.slide = slide

    @property
    def header(self) -> Section:
        """Slide的页眉区域

        Returns:
            Section: 返回一个Section对象
        """
        return Section(
            left=Cm(0), top=Cm(0), width=self.prs.slide_width, height=Cm(2.72)
        )

    @property
    def body(self) -> Section:
        """Slide的页面主体区域

        Returns:
            Section: 返回一个Section对象
        """
        return Section(
            left=Cm(0), top=Cm(2.91), width=self.prs.slide_width, height=Cm(13.81)
        )

    @property
    def footer(self) -> Section:
        """Slide的页脚区域

        Returns:
            Section: 返回一个Section对象
        """
        return Section(
            left=Cm(5.6), top=Cm(17.36), width=self.prs.slide_width, height=Cm(1.71)
        )

    def add_text(
        self,
        text: str,
        width: Union[Inches, Cm, int],
        height: Union[Inches, Cm, int],
        loc: Loc,
        anchor: Literal[
            "center",
            "top_left",
            "left_top",
            "top_right",
            "right_top",
            "top_mid",
            "mid_top",
            "top",
            "bottom_right",
            "right_bottom",
            "bottom_left",
            "left_bottom",
            "bottom_mid",
            "mid_bottom",
            "bottom",
            "mid_right",
            "right_mid",
            "right",
            "mid_left",
            "right_left",
            "left",
        ] = "center",
        shape_type: MSO_SHAPE_TYPE = MSO_SHAPE.RECTANGLE,
        fill_color: Optional[Union[RGBColor, str]] = None,
        border_color: Optional[Union[RGBColor, str]] = None,
        font_color: Optional[Union[RGBColor, str]] = None,
        font_size: Optional[Union[Pt, int]] = Pt(14),
        font_bold=False,
        font_italic=False,
        text_wrap: bool = False,
        *args,
        **kwargs,
    ) -> Shape:
        """_summary_

        Args:
            text (str): 文本内容
            width (Union[Inches, Cm, int]): 宽度
            height (Union[Inches, Cm, int]): 高度
            loc (Loc - left, top): _description_
            anchor (Literal[
                "center",
                "top_left", "left_top",
                "top_right", "right_top",
                "top_mid", "mid_top", "top",
                "bottom_right", "right_bottom",
                "bottom_left", "left_bottom",
                "bottom_mid", "mid_bottom", "bottom",
                "mid_right", "right_mid", "right",
                "mid_left", "right_left", "left",
            ]): 锚点. Defaults to "center".
            fill_color (Optional[Union[RGBColor, str]], optional): 填充颜色. Defaults to None.
            border_color (Optional[Union[RGBColor, str]], optional): 边框颜色. Defaults to None.
            font_color (Optional[Union[RGBColor, str]], optional): 字体颜色. Defaults to None.
            font_size (Optional[Union[Pt, int]], optional): 文字大小. Defaults to Pt(14).
            font_bold (bool, optional): 文字是否粗体. Defaults to False.
            font_italic (bool, optional): 文字是否斜体. Defaults to False.
            text_wrap (bool, optional): 文字是否换行. Defaults to False.

        Returns:
            Shape: 带有文本的形状
        """
        left, top = AnchorLoc(loc.left, loc.top, width, height, anchor=anchor).loc
        shape = self.slide.shapes.add_shape(
            shape_type,
            left=left,
            top=top,
            width=width,
            height=height,
        )

        # 形状填充
        fill = shape.fill
        if fill_color is None:
            fill.background()
        else:
            fill.solid()
            fill.fore_color.rgb = (
                fill_color
                if isinstance(fill_color, RGBColor)
                else RGBColor.from_string(fill_color)
            )

        # 边框
        line = shape.line
        if border_color is None:
            line.fill.background()
        else:
            line.color.rgb = (
                border_color
                if isinstance(border_color, RGBColor)
                else RGBColor.from_string(border_color)
            )

        # 形状文字
        text_frame = shape.text_frame
        text_frame.word_wrap = text_wrap
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text

        # 设置字体
        font = run.font
        font.size = font_size if isinstance(font_size, Pt) else Pt(font_size)
        font.bold = font_bold
        font.italic = font_italic
        if font_color is None:
            if is_light_color(fill_color):
                font.color.rgb = RGBColor(0, 0, 0)
            else:
                font.color.rgb = RGBColor(255, 255, 255)
        else:
            font.color.rgb = (
                font_color
                if isinstance(font_color, RGBColor)
                else RGBColor.from_string(font_color)
            )

        return shape

    def add_image(
        self,
        image_file: str,
        width: Optional[Union[Inches, Cm, int]] = None,
        height: Optional[Union[Inches, Cm, int]] = None,
        top: Optional[Union[Inches, Cm, int]] = None,
        left: Optional[Union[Inches, Cm, int]] = None,
        loc: Optional[Loc] = None,
        anchor: Optional[
            Literal[
                "center",
                "top_left",
                "left_top",
                "top_right",
                "right_top",
                "top_mid",
                "mid_top",
                "top",
                "bottom_right",
                "right_bottom",
                "bottom_left",
                "left_bottom",
                "bottom_mid",
                "mid_bottom",
                "bottom",
                "mid_right",
                "right_mid",
                "right",
                "mid_left",
                "right_left",
                "left",
            ]
        ] = "center",
    ) -> Picture:
        # 先插入图片，方便获取图片长宽（插入图片时如不同时指定宽和高，pptx高会自动缩放尺寸），再调整位置
        image = self.slide.shapes.add_picture(
            image_file=image_file,
            left=0 if left is None else left,
            top=0 if top is None else top,
            width=width,
            height=height,
        )

        if loc is not None:
            left, top = AnchorLoc(
                loc.left, loc.top, image.width, image.height, anchor=anchor
            ).loc
            image.left = left
            image.top = top

        return image

    def set_title(
        self, title: str, font_size: Optional[Union[Pt, Inches, Cm]] = None
    ) -> Shape:
        """设置幻灯片标题

        Args:
            title (str): 标题内容
            font_size (Optional[Union[Pt, Inches, Cm]], optional): 标题字体大小. Defaults to None.

        Returns:
            Shape: 带有标题的形状
        """

        title_shape = self.slide.shapes.title
        title_shape.text = title

        # 设置字体大小
        if font_size is not None:
            text_frame = title_shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = font_size

        return title_shape


class PPT:
    def __init__(self, template_path: str, save_path: str = None) -> None:
        self.template_path = template_path
        self.save_path = (
            save_path or f"{os.path.dirname(self.template_path)}presentation.pptx"
        )
        self.prs = Presentation(template_path)
        """初始化一个PPT类

        Args:
            template_path (str): 模板文件路径
            save_path (str): 保存路径，默认为和模板路径同一文件夹下的presentation.pptx文件. Defaults to None.
            
        Returns:
            _type_: _description_
        """

    @property
    def parent_slide(self) -> SlideLayout:
        """内容页幻灯片母版

        Returns
        -------
        SlideLayout
            返回该ppt的内容页幻灯片母版对象
        """
        slide = self.prs.slide_layouts[0]
        return slide

    def add_content_slide(
        self,
        layout_style: int = 0,
    ) -> SlideContent:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_style])
        print(f"已添加第{self.prs.slides.index(slide)}页.")
        content = SlideContent(self.prs, slide)
        return content

    def save(self):
        self.prs.save(self.save_path)
        print(f"{self.save_path} has been saved")


if __name__ == "__main__":
    p = PPT("template.pptx")
    c = p.add_content_slide()
    c.set_title("测试测试")
    c.add_text(
        "测试",
        width=Cm(4),
        height=Cm(0.5),
        loc=c.body.top_right,
        anchor="top_right",
        fill_color="F0CB46",
    )
    c.add_image(
        "D:\PyProjects\chart_class\plots\肾性贫血市场分治疗大类PTD滚动年趋势.png",
        width=c.body.width * 0.8,
        height=None,
        loc=c.body.center,
    )
    p.save()
